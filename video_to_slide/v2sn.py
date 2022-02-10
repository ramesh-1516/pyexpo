import os
import time
import cv2
import imutils
import shutil
import time
import pytesseract
from pptx import Presentation
from pptx.util import Inches
import tkinter as tk
from tkinter import *
import youtube_dl
from tkinter import messagebox, filedialog

OUTPUT_SLIDES_DIR = f"./output"
framerate = 3  # no.of frames per second that needs to be processed, fewer the count faster the speed
warmup = framerate  # initial number of frames to be skipped
fgbshistory = framerate * 15  # no.of frames in background object
varthreshold = 16  # Threshold on the squared Mahalanobis distance between the pixel and the model to decide whether a pixel is well described by the background model.
detectshadows = False  # If true, the algorithm will detect shadows and mark them.
minipercent = 0.1  # min % of diff between foreground and background to detect if motion has stopped
maxpercent = 3  # max % of diff between foreground and background to detect if frame is still in motion


def get_frames(video_path):
    vs = cv2.VideoCapture(video_path)
    framerate = vs.get(cv2.CAP_PROP_FPS)
    if not vs.isOpened():
        raise Exception(f'unable to open file {video_path}')

    total_frames = vs.get(cv2.CAP_PROP_FRAME_COUNT)
    frame_time = 0
    frame_count = 0
    print("total_frames: ", total_frames)
    print("framerate", framerate)

    while True:

        vs.set(cv2.CAP_PROP_POS_MSEC, frame_time * 1000)  # move frame to a timestamp
        frame_time += 1 / framerate

        (_, frame) = vs.read()

        if frame is None:
            break

        frame_count += 1
        yield frame_count, frame_time, frame

    vs.release()


pytesseract.pytesseract.tesseract_cmd = 'C:\Program Files\Tesseract-OCR/tesseract.exe'


def detect_unique_screenshots(video_path, output_folder_screenshot_path):
    fgbg = cv2.createBackgroundSubtractorMOG2(history=warmup, varThreshold=varthreshold, detectShadows=detectshadows)

    captured = False
    start_time = time.time()
    (W, H) = (None, None)

    screenshoots_count = 0
    for frame_count, frame_time, frame in get_frames(video_path):
        orig = frame.copy()  # clone the original frame (so we can save it later),
        frame = imutils.resize(frame, width=600)  # resize the frame
        mask = fgbg.apply(frame)  # apply the background subtractor

        if W is None or H is None:
            (H, W) = mask.shape[:2]

        p_diff = (cv2.countNonZero(mask) / float(W * H)) * 100

        if p_diff < minipercent and not captured and frame_count > warmup:
            captured = True
            filename = f"{screenshoots_count:03}_{round(frame_time / 60, 2)}.png"

            path = os.path.join(output_folder_screenshot_path, filename)
            print("saving {}".format(path))
            cv2.imwrite(path, orig)
            img1 = cv2.imread(path)
            img2 = cv2.cvtColor(img1, cv2.COLOR_BGR2RGB)
            convert_screenshots_to_pptx(img2, path)
            screenshoots_count += 1

        elif captured and p_diff >= maxpercent:
            captured = False
    print(f'{screenshoots_count} screenshots Captured!')
    print(f'Time taken {time.time() - start_time}s')
    return


def initialize_output_folder(video_path):
    '''Clean the output folder if already exists'''
    output_folder_screenshot_path = f"{OUTPUT_SLIDES_DIR}/{video_path.rsplit('/')[-1].split('.')[0]}"

    if os.path.exists(output_folder_screenshot_path):
        shutil.rmtree(output_folder_screenshot_path)

    os.makedirs(output_folder_screenshot_path, exist_ok=True)
    print('initialized output folder', output_folder_screenshot_path)
    return output_folder_screenshot_path


prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)


def convert_screenshots_to_pptx(img2, path):
    mytext = ""
    i = 0
    k = []
    texts = pytesseract.image_to_data(img2)
    img_path: str = path
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    top = Inches(1)
    left1 = Inches(12)
    height1 = width1 = Inches(3.5)
    for b, o in enumerate(texts.splitlines()):
        g = o.split()
        k.append(g)
    g.append("")
    g[6] = 0
    k.append(g)
    m = 10000
    lst1 = [1]
    lst2 = [1]
    for j in k:
        if len(j) == 12 and i >= 1:
            x, y, w, h = int(j[6]), int(j[7]), int(j[8]), int(j[9])
            if x < m and i >= 2:
                txBox = slide.shapes.add_textbox(left=7500 * min(lst1), top=6000 * (sum(lst2) / len(lst2)), width=w,
                                                 height=h)
                tf = txBox.text_frame
                tf.text = mytext
                # time.sleep(1)
                mytext = ""
                lst1.clear()
                lst2.clear()
            mytext += j[11] + " "
            m = x
            lst1.append(x)
            lst2.append(y)
        i += 1
    pic = slide.shapes.add_picture(img_path, left=left1, top=top, height=height1, width=width1)
    prs.save('output/hellohi.pptx')


if __name__ == "__main__":

    def Widgets():
        head_label = Label(root, text=" select video from PC or from youtube ", bg="#34abeb", fg="#0003c9")
        head_label.grid(row=0, column=1, pady=10, padx=5)

        pcBrowse_B = Button(root, text=" Video from PC ", command=close, width=20, bg="thistle1", pady=10, padx=15,
                            relief=GROOVE, font="Georgia, 13")
        pcBrowse_B.grid(row=1, column=1, pady=5, padx=5)

        link_label = Label(root, text="YouTube link :", bg="#34eba5", pady=5, padx=5)
        link_label.grid(row=2, column=0, pady=5, padx=5)

        root.linkText = Entry(root, width=35, textvariable=video_Link, font="Arial 14")
        root.linkText.grid(row=2, column=1, pady=5, padx=5, columnspan=2)

        destination_label = Label(root, text="Destination :", bg="#34eba5", pady=5, padx=9)
        destination_label.grid(row=3, column=0, pady=5, padx=5)

        root.destinationText = Entry(root, width=27, textvariable=download_Path, font="Arial 14")
        root.destinationText.grid(row=3, column=1, pady=5, padx=5)

        browse_B = Button(root, text="Browse", command=Browse, width=10, bg="bisque", relief=GROOVE)
        browse_B.grid(row=3, column=2, pady=1, padx=1)

        Download_B = Button(root, text="Download Video", command=Download, width=20, bg="thistle1", pady=10, padx=15,
                            relief=GROOVE, font="Georgia, 13")
        Download_B.grid(row=4, column=1, pady=5, padx=5)


    def close():
        root.destroy()


    def Browse():
        download_Directory = filedialog.askdirectory(initialdir="YOUR DIRECTORY PATH", title="Save Video")
        download_Path.set(download_Directory)


    def Download():

        link_of_the_video = video_Link.get()
        download_Folder = download_Path.get()
        ydl_opts = {
            'format': 'best',
            'outtmpl': download_Folder + '/%(title)s.%(ext)s',
        }
        yt = link_of_the_video.strip()
        with youtube_dl.YoutubeDL(ydl_opts) as ydl:
            ydl.download([yt])
        messagebox.showinfo("SUCCESSFULLY", "DOWNLOADED AND SAVED IN\n" + download_Folder)
        root.destroy()


    def openFile():
        filepath = filedialog.askopenfilename()
        return filepath


    check = "y"
    while check == "y":
        root = tk.Tk()
        root.geometry("530x280")
        root.resizable(False, False)
        root.title("select video from PC or from youtube")
        root.config(background="#34abeb")
        video_Link = StringVar()
        download_Path = StringVar()
        Widgets()
        root.mainloop()
        video_path = openFile()

        print('video_path', video_path)
        output_folder_screenshot_path = initialize_output_folder(video_path)
        detect_unique_screenshots(video_path, output_folder_screenshot_path)
        check = input("If you want to add another video (y/n) :")